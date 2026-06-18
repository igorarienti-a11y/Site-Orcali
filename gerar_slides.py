from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Cores
BG          = RGBColor(0x0B, 0x1E, 0x1F)
BG_GRAY     = RGBColor(0xE8, 0xE8, 0xE8)
GREEN       = RGBColor(0x39, 0xE5, 0x5F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1A, 0x2A, 0x2A)
GRAY        = RGBColor(0xB0, 0xC4, 0xC4)
GOOGLE_BLUE = RGBColor(0x42, 0x85, 0xF4)
META_BLUE   = RGBColor(0x00, 0x82, 0xFB)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, radius=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    if radius:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape


def add_text(slide, text, l, t, w, h, font_size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, l, t, w, h, font_size=13, color=WHITE, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_green_line(slide, l, t, w, h=0.04):
    return add_rect(slide, l, t, w, h, fill_color=GREEN)


def add_icon(slide, l, t, size=0.55):
    add_rect(slide, l, t, size, size, fill_color=RGBColor(0x1A, 0x30, 0x30), radius=20000)
    txt = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(size), Inches(size))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'N'
    run.font.size = Pt(18)
    run.font.color.rgb = GREEN
    run.font.bold = True


def pill_badge(slide, text, l, t):
    badge = add_rect(slide, l, t, 2.0, 0.38,
                     fill_color=RGBColor(0x1A, 0x30, 0x30),
                     line_color=RGBColor(0x2A, 0x4A, 0x4A),
                     radius=50000)
    txt = slide.shapes.add_textbox(Inches(l + 0.1), Inches(t + 0.05), Inches(1.8), Inches(0.28))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'N  ' + text
    run.font.size = Pt(13)
    run.font.color.rgb = GREEN
    run.font.bold = True


# ─────────────────────────────────────────
# SLIDE 1 — DIVIDER
# ─────────────────────────────────────────
s1 = blank_slide(prs)
set_bg(s1, BG)
pill_badge(s1, 'Planejamento Digital', 5.55, 0.85)
add_text(s1, 'Performance &', 2.5, 1.7, 8.33, 1.2,
         font_size=58, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, 'Aquisicao de Leads', 2.5, 2.75, 8.33, 1.3,
         font_size=58, color=GREEN, align=PP_ALIGN.CENTER)
add_green_line(s1, 5.9, 4.25, 1.5)
add_text(s1, 'Trafego pago  |  Landing Pages  |  Integracoes',
         2.5, 4.5, 8.33, 0.5, font_size=15, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 2 — OBJETIVOS (split)
# ─────────────────────────────────────────
s2 = blank_slide(prs)
set_bg(s2, BG_GRAY)
add_rect(s2, 4.8, 0, 8.53, 7.5, fill_color=BG)
add_icon(s2, 0.35, 0.35)
add_text(s2, 'Objetivos', 0.35, 1.5, 4.0, 1.2,
         font_size=46, color=DARK)
add_text(s2, 'de performance', 0.35, 2.85, 4.0, 0.45,
         font_size=14, color=RGBColor(0x5A, 0x7A, 0x7A))
items = [
    ('1', 'Alcance e reconhecimento de marca'),
    ('2', 'Geracao de leads qualificados'),
    ('3', 'Mensuracao de cada lead gerado'),
]
for i, (num, label) in enumerate(items):
    y = 1.8 + i * 1.4
    add_rect(s2, 5.5, y, 0.5, 0.5,
             fill_color=GREEN, radius=8000)
    add_text(s2, num, 5.5, y, 0.5, 0.5,
             font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s2, label, 6.2, y + 0.05, 5.8, 0.5,
             font_size=18, color=WHITE)

# ─────────────────────────────────────────
# SLIDE 3 — VERBA R$8.000
# ─────────────────────────────────────────
s3 = blank_slide(prs)
set_bg(s3, BG)
add_icon(s3, 0.35, 0.35)
add_text(s3, 'Investimento Mensal', 1.1, 0.38, 6.0, 0.4,
         font_size=12, color=GREEN)
add_text(s3, 'R$ 8.000 / mes', 0.55, 0.9, 8.0, 1.0,
         font_size=44, color=WHITE)
add_green_line(s3, 0.55, 2.0, 12.2)
# Card Google
add_rect(s3, 0.55, 2.2, 5.9, 4.7, fill_color=RGBColor(0x0F, 0x26, 0x28))
add_rect(s3, 0.55, 2.2, 5.9, 0.07, fill_color=GOOGLE_BLUE)
add_text(s3, 'GOOGLE ADS', 0.85, 2.38, 3.0, 0.4,
         font_size=11, color=GRAY)
add_text(s3, 'R$ 5.500', 0.85, 2.8, 4.0, 0.8,
         font_size=38, color=WHITE)
add_text(s3, '69% do investimento total', 0.85, 3.55, 4.0, 0.35,
         font_size=13, color=GRAY)
add_multiline(s3,
    ['Search Facilities  -  R$ 2.200',
     'Search Seguranca  -  R$ 1.800',
     'Display Topo de Funil  -  R$ 1.000'],
    0.85, 4.05, 5.2, 1.6,
    font_size=13, color=RGBColor(0xC0, 0xD4, 0xD4))
# Card Meta
add_rect(s3, 6.88, 2.2, 5.9, 4.7, fill_color=RGBColor(0x0F, 0x26, 0x28))
add_rect(s3, 6.88, 2.2, 5.9, 0.07, fill_color=META_BLUE)
add_text(s3, 'META ADS', 7.18, 2.38, 3.0, 0.4,
         font_size=11, color=GRAY)
add_text(s3, 'R$ 2.500', 7.18, 2.8, 4.0, 0.8,
         font_size=38, color=WHITE)
add_text(s3, '31% do investimento total', 7.18, 3.55, 4.0, 0.35,
         font_size=13, color=GRAY)
add_multiline(s3,
    ['Alcance Institucional  -  R$ 1.600',
     'Retargeting LP  -  R$ 900'],
    7.18, 4.05, 5.2, 1.2,
    font_size=13, color=RGBColor(0xC0, 0xD4, 0xD4))

# ─────────────────────────────────────────
# SLIDE 4 — META: ALCANCE
# ─────────────────────────────────────────
s4 = blank_slide(prs)
set_bg(s4, BG)
add_icon(s4, 0.35, 0.35)
add_text(s4, 'Meta Ads  |  R$ 1.600/mes  |  Objetivo: Reach',
         1.1, 0.38, 10.0, 0.4, font_size=12, color=GREEN)
add_text(s4, 'Alcance institucional por nicho',
         0.55, 0.9, 9.0, 0.8, font_size=36, color=WHITE)
add_green_line(s4, 0.55, 1.82, 12.2)
nicho_data = [
    ('Saude',
     '"Operacoes hospitalares nao podem parar.\nPor tras de cada turno funcionando,\nexiste uma estrutura invisivel."'),
    ('Ensino',
     '"De manha cedo ate o fim da noite,\num campus nunca para.\nA gestao que garante isso, sim."'),
    ('Financeiro',
     '"Compliance comeca antes da portaria.\nFacilities e seguranca integradas\npara instituicoes financeiras."'),
]
for i, (title, quote) in enumerate(nicho_data):
    x = 0.55 + i * 4.3
    add_rect(s4, x, 2.05, 4.0, 5.0, fill_color=RGBColor(0x0F, 0x28, 0x28))
    add_green_line(s4, x, 2.05, 4.0, 0.06)
    add_text(s4, title, x + 0.3, 2.25, 3.4, 0.5,
             font_size=20, bold=True, color=GREEN)
    add_text(s4, quote, x + 0.3, 2.9, 3.5, 3.8,
             font_size=13, color=WHITE)
add_text(s4,
    'Publico: Lookalike 1% SC/PR a partir de 2.966 clientes ativos',
    0.55, 7.1, 12.0, 0.35, font_size=12, color=GRAY)

# ─────────────────────────────────────────
# SLIDE 5 — META: RETARGETING
# ─────────────────────────────────────────
s5 = blank_slide(prs)
set_bg(s5, BG)
add_icon(s5, 0.35, 0.35)
add_text(s5, 'Meta Ads  |  R$ 900/mes',
         1.1, 0.38, 8.0, 0.4, font_size=12, color=GREEN)
add_text(s5, 'Retargeting',
         0.55, 0.9, 8.0, 0.8, font_size=36, color=WHITE)
add_text(s5, 'Convertendo quem ja demonstrou interesse',
         0.55, 1.6, 9.0, 0.45, font_size=16, color=GRAY)
add_green_line(s5, 0.55, 2.15, 12.2)
steps = [
    ('Viu o alcance',     'Engajou no Instagram\nou visitou o perfil'),
    ('Acessou a LP',      'Visitante que nao preencheu\no formulario'),
    ('Recebe o anuncio',  'Prova social + CTA direto\npara o formulario'),
    ('Preenche o form',   'Lead qualificado\nentra no funil'),
]
for i, (title, desc) in enumerate(steps):
    x = 0.55 + i * 3.2
    fc = RGBColor(0x0D, 0x24, 0x1A) if i == 2 else RGBColor(0x0F, 0x28, 0x28)
    add_rect(s5, x, 2.38, 2.9, 3.2, fill_color=fc)
    if i == 2:
        add_green_line(s5, x, 2.38, 2.9, 0.06)
    c = GREEN if i == 2 else WHITE
    add_text(s5, title, x + 0.2, 2.65, 2.5, 0.55,
             font_size=15, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s5, desc, x + 0.2, 3.3, 2.5, 1.5,
             font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s5, '>', x + 2.9, 3.4, 0.3, 0.4,
                 font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s5, 0.55, 5.85, 12.2, 1.3, fill_color=RGBColor(0x0F, 0x28, 0x28))
add_multiline(s5,
    ['Janela: 30 dias de visita a LP',
     'Angulo: "1.554 clientes. Contrato medio de 20 anos."',
     'CTA: Solicitar diagnostico operacional gratuito - resposta em 24h'],
    0.85, 5.95, 11.5, 1.1,
    font_size=13, color=RGBColor(0xC0, 0xD4, 0xD4))

# ─────────────────────────────────────────
# SLIDE 6 — GOOGLE: SEARCH
# ─────────────────────────────────────────
s6 = blank_slide(prs)
set_bg(s6, BG)
add_icon(s6, 0.35, 0.35)
add_text(s6, 'Google Ads  |  R$ 4.000/mes  |  Search',
         1.1, 0.38, 9.0, 0.4, font_size=12, color=GOOGLE_BLUE)
add_text(s6, 'Capturar quem ja busca pelo servico',
         0.55, 0.9, 10.0, 0.8, font_size=36, color=WHITE)
add_green_line(s6, 0.55, 1.82, 12.2)
# Card Facilities
add_rect(s6, 0.55, 2.05, 6.0, 5.1, fill_color=RGBColor(0x0F, 0x26, 0x28))
add_rect(s6, 0.55, 2.05, 6.0, 0.07, fill_color=GOOGLE_BLUE)
add_text(s6, 'Facilities', 0.85, 2.22, 3.5, 0.5,
         font_size=20, bold=True, color=WHITE)
add_text(s6, 'R$ 2.200/mes', 4.2, 2.22, 2.0, 0.5,
         font_size=15, color=WHITE, align=PP_ALIGN.RIGHT)
kws_f = [
    'empresa de facilities SC',
    'terceirizacao portaria Floripa',
    'gestao predial hospital',
    'limpeza terceirizada universidade',
    'facilities Curitiba',
    'zeladoria escola SC',
    'facilities instituicao financeira',
]
for i, kw in enumerate(kws_f):
    col = i % 2
    row = i // 2
    add_rect(s6, 0.75 + col * 2.9, 2.9 + row * 0.75, 2.7, 0.55,
             fill_color=RGBColor(0x14, 0x32, 0x38), radius=6000)
    add_text(s6, kw, 0.85 + col * 2.9, 2.95 + row * 0.75, 2.5, 0.42,
             font_size=11, color=RGBColor(0xA0, 0xC8, 0xFF))
# Card Segurança
add_rect(s6, 6.88, 2.05, 6.0, 5.1, fill_color=RGBColor(0x0F, 0x26, 0x28))
add_rect(s6, 6.88, 2.05, 6.0, 0.07, fill_color=RGBColor(0x34, 0xA8, 0x53))
add_text(s6, 'Seguranca', 7.18, 2.22, 3.5, 0.5,
         font_size=20, bold=True, color=WHITE)
add_text(s6, 'R$ 1.800/mes', 10.6, 2.22, 2.0, 0.5,
         font_size=15, color=WHITE, align=PP_ALIGN.RIGHT)
kws_s = [
    'vigilancia patrimonial SC',
    'empresa de seguranca Floripa',
    'seguranca hospitalar SC',
    'vigilante terceirizado Curitiba',
    'seguranca escolar Joinville',
    'seguranca instituicao financeira',
]
for i, kw in enumerate(kws_s):
    col = i % 2
    row = i // 2
    add_rect(s6, 7.05 + col * 2.9, 2.9 + row * 0.75, 2.7, 0.55,
             fill_color=RGBColor(0x0F, 0x28, 0x1E), radius=6000)
    add_text(s6, kw, 7.15 + col * 2.9, 2.95 + row * 0.75, 2.5, 0.42,
             font_size=11, color=RGBColor(0x7D, 0xCF, 0xA0))
add_text(s6,
    'Ad groups segmentados por nicho: copy e extensoes personalizadas para Saude, Ensino e Financeiro',
    0.55, 7.1, 12.2, 0.35, font_size=12, color=GRAY)

# ─────────────────────────────────────────
# SLIDE 7 — GOOGLE: DISPLAY
# ─────────────────────────────────────────
s7 = blank_slide(prs)
set_bg(s7, BG)
add_icon(s7, 0.35, 0.35)
add_text(s7, 'Google Ads  |  R$ 1.000/mes  |  Display',
         1.1, 0.38, 9.0, 0.4, font_size=12, color=GOOGLE_BLUE)
add_text(s7, 'Alcancar decisores antes da busca ativa',
         0.55, 0.9, 10.0, 0.8, font_size=36, color=WHITE)
add_green_line(s7, 0.55, 1.82, 12.2)
bloco_data = [
    ('Segmentacao por cargo',
     'Diretores Administrativos, Gerentes de Facilities e Gestores Operacionais em SC e PR.',
     'In-Market'),
    ('Segmentacao por empresa',
     'Hospitais, universidades, cooperativas e bancos na rede Google + YouTube.',
     'Audience'),
    ('Alimenta o retargeting',
     'Quem clica no Display entra na LP e cai na audiencia de retargeting do Meta. Funil integrado.',
     'Funil'),
]
for i, (title, desc, tag) in enumerate(bloco_data):
    x = 0.55 + i * 4.3
    add_rect(s7, x, 2.05, 4.0, 3.7, fill_color=RGBColor(0x0F, 0x26, 0x28))
    add_text(s7, title, x + 0.3, 2.2, 3.4, 0.6,
             font_size=16, bold=True, color=WHITE)
    add_text(s7, desc, x + 0.3, 2.85, 3.4, 2.0,
             font_size=13, color=GRAY)
    add_rect(s7, x + 0.3, 5.4, 1.2, 0.32,
             fill_color=RGBColor(0x14, 0x32, 0x46), radius=6000)
    add_text(s7, tag, x + 0.3, 5.42, 1.2, 0.28,
             font_size=10, color=RGBColor(0xA0, 0xC8, 0xFF), align=PP_ALIGN.CENTER)
add_rect(s7, 0.55, 6.1, 12.2, 1.05, fill_color=RGBColor(0x0F, 0x26, 0x28))
add_text(s7,
    'Display e o gerador de audiencia para o retargeting - nao e onde o lead converte, e onde ele descobre a Orcali.',
    0.85, 6.25, 11.0, 0.75, font_size=14, color=WHITE)

# ─────────────────────────────────────────
# SLIDE 8 — LP + TYPEBOT
# ─────────────────────────────────────────
s8 = blank_slide(prs)
set_bg(s8, BG)
add_icon(s8, 0.35, 0.35)
add_text(s8, 'Estrutura de Captacao',
         1.1, 0.38, 8.0, 0.4, font_size=12, color=GREEN)
add_text(s8, 'Landing Page + Formulario de qualificacao',
         0.55, 0.9, 11.0, 0.8, font_size=34, color=WHITE)
add_green_line(s8, 0.55, 1.82, 12.2)
# LP
add_rect(s8, 0.55, 2.05, 6.0, 5.1, fill_color=RGBColor(0x0F, 0x28, 0x28))
add_rect(s8, 0.55, 2.05, 6.0, 0.55, fill_color=RGBColor(0x00, 0x2E, 0x6E))
add_text(s8, 'orcali.com.br/facilities',
         0.85, 2.1, 3.8, 0.45, font_size=12, color=RGBColor(0xC0, 0xD4, 0xD4))
add_rect(s8, 5.3, 2.1, 0.95, 0.35,
         fill_color=GREEN, radius=6000)
add_text(s8, 'PRONTO', 5.3, 2.13, 0.95, 0.3,
         font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
lp_items = [
    'Hero: "Operacoes que nao podem parar"',
    'Diferenciais: 58 anos | 1.554 clientes',
    'Servicos por nicho: Saude, Ensino, Financeiro',
    'Prova social + certificacoes ISO',
    'Formulario integrado via Typebot',
    'LP Seguranca  ->  Mes 2',
]
for i, item in enumerate(lp_items):
    c = GRAY if i == 5 else WHITE
    add_text(s8, '->  ' + item,
             0.75, 2.82 + i * 0.38, 5.6, 0.38,
             font_size=13, color=c)
# Form
add_rect(s8, 6.88, 2.05, 6.0, 5.1, fill_color=RGBColor(0x0F, 0x28, 0x28))
add_text(s8, 'Campos do formulario',
         7.18, 2.15, 5.0, 0.4, font_size=13, bold=True, color=GREEN)
fields = [
    'Nome completo',
    'Empresa',
    'CNPJ',
    'E-mail corporativo',
    'Telefone / WhatsApp',
    'Segmento  (Saude | Ensino | Financeiro)',
    'Cidade / Regiao',
]
for i, field in enumerate(fields):
    add_rect(s8, 7.08, 2.7 + i * 0.43, 5.5, 0.36,
             fill_color=RGBColor(0x16, 0x30, 0x30), radius=4000)
    add_text(s8, '  ' + field,
             7.22, 2.72 + i * 0.43, 5.2, 0.32,
             font_size=12, color=WHITE)
add_rect(s8, 7.08, 5.82, 5.5, 0.5,
         fill_color=GREEN, radius=6000)
add_text(s8, 'Solicitar diagnostico operacional gratuito  ->',
         7.08, 5.84, 5.5, 0.46,
         font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 9 — INTEGRACOES
# ─────────────────────────────────────────
s9 = blank_slide(prs)
set_bg(s9, BG)
add_icon(s9, 0.35, 0.35)
add_text(s9, 'Integracoes & Mensuracao',
         1.1, 0.38, 8.0, 0.4, font_size=12, color=GREEN)
add_text(s9, 'Do clique ao comercial rastreado',
         0.55, 0.9, 10.0, 0.8, font_size=36, color=WHITE)
add_green_line(s9, 0.55, 1.82, 12.2)
nodes = [
    ('Meta / Google',    'Pixel + CAPI\nGoogle Tag',    META_BLUE),
    ('Landing Page',     'Pixel fire\nViewContent',     GREEN),
    ('Planilha Google',  'Lead + UTMs\n+ origem',       RGBColor(0x0F, 0x9D, 0x58)),
    ('RD Station CRM',   'Nutricao\n+ funil',           RGBColor(0x8B, 0x5C, 0xF6)),
    ('Comercial',        'Lead pronto\np/ abordagem',   RGBColor(0xFB, 0x92, 0x3C)),
]
for i, (title, desc, color) in enumerate(nodes):
    x = 0.55 + i * 2.55
    add_rect(s9, x, 2.2, 2.15, 2.8, fill_color=RGBColor(0x0F, 0x28, 0x28))
    add_rect(s9, x, 2.2, 2.15, 0.07, fill_color=color)
    add_text(s9, title, x + 0.15, 2.4, 1.85, 0.6,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s9, desc, x + 0.15, 3.1, 1.85, 0.9,
             font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s9, '>', x + 2.15, 3.3, 0.4, 0.4,
                 font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
labels = ['clique', 'form submit', 'Apps Script', 'alerta CRM']
for i, label in enumerate(labels):
    x = 0.55 + i * 2.55 + 2.15
    add_text(s9, label, x, 5.1, 0.55, 0.3,
             font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s9, 0.55, 5.55, 12.2, 1.6, fill_color=RGBColor(0x0F, 0x28, 0x28))
add_text(s9,
    'Cada lead chega com origem rastreada: plataforma | campanha | nicho | palavra-chave.',
    0.85, 5.65, 11.0, 0.5, font_size=14, bold=True, color=WHITE)
add_text(s9,
    'Resolve a principal dor do historico: R$ 6k/mes em 2024 sem mensuracao de resultados.',
    0.85, 6.2, 11.0, 0.5, font_size=13, color=GRAY)

# ─────────────────────────────────────────
prs.save('C:/Users/igori/projetos/orcali-lp/orcali-trafego-slides.pptx')
print('Salvo: orcali-trafego-slides.pptx')
