"""
Gera slides de trafego clonando o design real do Orcali.pptx.
Clona slides existentes (preservando backgrounds, icones, gradientes)
e substitui apenas o conteudo textual.
"""
import copy, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC  = r'C:/Users/igori/Downloads/Orcali.pptx'
DEST = r'C:/Users/igori/projetos/orcali-lp/orcali-trafego-v2.pptx'

prs = Presentation(SRC)
W = prs.slide_width.inches    # 20.0
H = prs.slide_height.inches   # 11.25

GREEN       = RGBColor(0x39, 0xE5, 0x5F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x0B, 0x1E, 0x1F)
GRAY        = RGBColor(0xB0, 0xC4, 0xC4)
LIGHT_GRAY  = RGBColor(0xD0, 0xE0, 0xE0)
GOOGLE_BLUE = RGBColor(0x42, 0x85, 0xF4)
META_BLUE   = RGBColor(0x00, 0x82, 0xFB)
CARD_BG     = RGBColor(0x0F, 0x28, 0x28)


# ─────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────

def clone_slide(prs, slide_idx):
    """Clona slide[slide_idx] e adiciona ao final da apresentacao."""
    src = prs.slides[slide_idx]
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)

    # Limpa spTree do novo slide
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)

    # Copia shapes do slide original
    for el in src.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))

    # Copia background
    src_cSld = src._element.find(qn('p:cSld'))
    new_cSld = new_slide._element.find(qn('p:cSld'))
    src_bg = src_cSld.find(qn('p:bg'))
    if src_bg is not None:
        new_bg = new_cSld.find(qn('p:bg'))
        if new_bg is not None:
            new_cSld.remove(new_bg)
        sp_tree_el = new_cSld.find(qn('p:spTree'))
        idx = list(new_cSld).index(sp_tree_el)
        new_cSld.insert(idx, copy.deepcopy(src_bg))

    # Copia relacionamentos de imagens do slide original
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    for rel in src.part.rels.values():
        if 'image' in rel.reltype or 'slide' in rel.reltype:
            try:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                # Atualiza referencias no XML
                for el in new_slide.shapes._spTree.iter():
                    for attr in [f'{{{ns_r}}}embed', f'{{{ns_r}}}link']:
                        if el.get(attr) == rel.rId:
                            el.set(attr, new_rId)
            except Exception:
                pass

    return new_slide


def remove_textboxes(slide, keep_names=None):
    """Remove todas as TextBoxes do slide. Se keep_names, mantem as listadas."""
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            if keep_names and shape.name in keep_names:
                continue
            to_remove.append(shape._element)
    for el in to_remove:
        sp_tree.remove(el)


def set_text(shape, text, font_size=None, bold=None, color=None, align=None):
    """Substitui texto de uma shape existente preservando a formatacao base."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_tb(slide, text, l, t, w, h,
           font_size=16, bold=False, color=WHITE,
           align=PP_ALIGN.LEFT, font_name='Calibri'):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx


def add_ml(slide, lines, l, t, w, h,
           font_size=15, color=WHITE, font_name='Calibri', spacing=3):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tx


def add_rect(slide, l, t, w, h, fill=None, line=None, radius=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    if radius:
        sp = s._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return s


def hline(slide, l, t, w, h=0.06, color=GREEN):
    add_rect(slide, l, t, w, h, fill=color)


# ─────────────────────────────────────────
# SLIDE 19 — DIVIDER (clone slide 1)
# ─────────────────────────────────────────
s = clone_slide(prs, 0)

# Encontra e edita TextBox 11 (titulo grande)
for shape in s.shapes:
    if shape.name == 'TextBox 11' and shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        # Linha 1: "Performance &"
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = 'Performance &'
        r1.font.color.rgb = WHITE
        r1.font.size = Pt(72)
        # Linha 2: "Aquisicao de Leads" em verde
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = 'Aquisicao de Leads'
        r2.font.color.rgb = GREEN
        r2.font.size = Pt(72)

    if shape.name == 'TextBox 12' and shape.has_text_frame:
        set_text(shape, 'Trafego pago  |  Landing Pages  |  Integracoes',
                 font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

print('Slide 19 (Divider) ok')

# ─────────────────────────────────────────
# SLIDE 20 — OBJETIVOS (clone slide 2)
# ─────────────────────────────────────────
s = clone_slide(prs, 1)

for shape in s.shapes:
    if shape.name == 'TextBox 14' and shape.has_text_frame:
        set_text(shape, 'de performance', font_size=14, color=RGBColor(0x5A,0x7A,0x7A))
    if shape.name == 'TextBox 15' and shape.has_text_frame:
        set_text(shape, 'Alcance e reconhecimento de marca', font_size=22, color=WHITE)
    if shape.name == 'TextBox 41' and shape.has_text_frame:
        set_text(shape, 'Geracao de leads qualificados por nicho', font_size=22, color=WHITE)
    if shape.name == 'TextBox 40' and shape.has_text_frame:
        set_text(shape, 'Mensuracao e origem de cada lead gerado', font_size=22, color=WHITE)

print('Slide 20 (Objetivos) ok')

# ─────────────────────────────────────────
# SLIDE 21 — VERBA R$8.000 (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

# Titulo
add_tb(s, 'Investimento Mensal', 1.12, 0.80, 6.0, 0.55, font_size=14, color=GREEN)
add_tb(s, 'Distribuicao de Verba', 1.12, 1.4, 12.0, 1.0, font_size=44, color=WHITE)
hline(s, 1.12, 2.55, 17.5)

# Card Google
add_rect(s, 1.12, 2.8, 8.6, 7.8, fill=CARD_BG)
add_rect(s, 1.12, 2.8, 8.6, 0.1, fill=GOOGLE_BLUE)
add_tb(s, 'GOOGLE ADS', 1.5, 3.0, 5.0, 0.5, font_size=13, color=GRAY)
add_tb(s, 'R$ 5.500', 1.5, 3.55, 6.0, 1.1, font_size=54, color=WHITE)
add_tb(s, '69% do investimento total', 1.5, 4.7, 6.0, 0.45, font_size=15, color=GRAY)
add_ml(s, [
    'Search Facilities  -  R$ 2.200',
    'Search Seguranca  -  R$ 1.800',
    'Display Topo de Funil  -  R$ 1.000',
], 1.5, 5.4, 7.5, 2.5, font_size=16, color=LIGHT_GRAY)

# Card Meta
add_rect(s, 10.28, 2.8, 8.6, 7.8, fill=CARD_BG)
add_rect(s, 10.28, 2.8, 8.6, 0.1, fill=META_BLUE)
add_tb(s, 'META ADS', 10.7, 3.0, 5.0, 0.5, font_size=13, color=GRAY)
add_tb(s, 'R$ 2.500', 10.7, 3.55, 6.0, 1.1, font_size=54, color=WHITE)
add_tb(s, '31% do investimento total', 10.7, 4.7, 6.0, 0.45, font_size=15, color=GRAY)
add_ml(s, [
    'Alcance Institucional  -  R$ 1.600',
    'Retargeting LP  -  R$ 900',
], 10.7, 5.4, 7.5, 1.8, font_size=16, color=LIGHT_GRAY)

print('Slide 21 (Verba) ok')

# ─────────────────────────────────────────
# SLIDE 22 — META: ALCANCE (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

add_tb(s, 'Meta Ads  |  R$ 1.600/mes  |  Objetivo: Reach',
       1.12, 0.80, 14.0, 0.55, font_size=14, color=GREEN)
add_tb(s, 'Alcance institucional por nicho',
       1.12, 1.4, 14.0, 1.0, font_size=44, color=WHITE)
hline(s, 1.12, 2.55, 17.5)

nichos = [
    ('Saude', '"Operacoes hospitalares nao podem parar.\nPor tras de cada turno funcionando,\nexiste uma estrutura invisivel."'),
    ('Ensino', '"De manha cedo ate o fim da noite,\num campus nunca para.\nA gestao que garante isso, sim."'),
    ('Financeiro', '"Compliance comeca antes da portaria.\nFacilities e seguranca integradas\npara instituicoes financeiras."'),
]
for i, (title, quote) in enumerate(nichos):
    x = 1.12 + i * 6.2
    add_rect(s, x, 2.8, 5.8, 7.8, fill=CARD_BG)
    hline(s, x, 2.8, 5.8, 0.1)
    add_tb(s, title, x+0.4, 3.1, 5.0, 0.75, font_size=24, bold=True, color=GREEN)
    add_tb(s, quote, x+0.4, 4.0, 5.0, 5.5, font_size=16, color=WHITE)

add_tb(s, 'Publico: Lookalike 1% SC/PR a partir de 2.966 clientes ativos',
       1.12, 10.7, 17.5, 0.45, font_size=13, color=GRAY)

print('Slide 22 (Meta Alcance) ok')

# ─────────────────────────────────────────
# SLIDE 23 — META: RETARGETING (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

add_tb(s, 'Meta Ads  |  R$ 900/mes',
       1.12, 0.80, 10.0, 0.55, font_size=14, color=GREEN)
add_tb(s, 'Retargeting',
       1.12, 1.4, 10.0, 1.0, font_size=44, color=WHITE)
add_tb(s, 'Convertendo quem ja demonstrou interesse',
       1.12, 2.45, 14.0, 0.6, font_size=20, color=GRAY)
hline(s, 1.12, 3.15, 17.5)

steps = [
    ('Viu o alcance',     'Engajou no Instagram\nou visitou o perfil'),
    ('Acessou a LP',      'Visitante que nao preencheu\no formulario'),
    ('Recebe o anuncio',  'Prova social + CTA direto'),
    ('Preenche o form',   'Lead qualificado\nentra no funil'),
]
for i, (title, desc) in enumerate(steps):
    x = 1.12 + i * 4.6
    is_highlight = (i == 2)
    fc = RGBColor(0x0D, 0x24, 0x1A) if is_highlight else CARD_BG
    add_rect(s, x, 3.4, 4.1, 4.5, fill=fc)
    if is_highlight:
        hline(s, x, 3.4, 4.1, 0.1)
    c = GREEN if is_highlight else WHITE
    add_tb(s, title, x+0.3, 3.8, 3.5, 0.7, font_size=18, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_tb(s, desc, x+0.3, 4.65, 3.5, 2.0, font_size=15, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_tb(s, '>', x+4.1, 5.2, 0.5, 0.6, font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s, 1.12, 8.3, 17.5, 2.1, fill=CARD_BG)
add_ml(s, [
    'Janela: 30 dias de visita a LP',
    'Angulo: "1.554 clientes. Contrato medio de 20 anos."',
    'CTA: Solicitar diagnostico operacional gratuito - resposta em 24h',
], 1.5, 8.55, 16.5, 1.7, font_size=15, color=LIGHT_GRAY)

print('Slide 23 (Meta Retargeting) ok')

# ─────────────────────────────────────────
# SLIDE 24 — GOOGLE: SEARCH (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

add_tb(s, 'Google Ads  |  R$ 4.000/mes  |  Search',
       1.12, 0.80, 12.0, 0.55, font_size=14, color=GOOGLE_BLUE)
add_tb(s, 'Capturar quem ja busca pelo servico',
       1.12, 1.4, 14.0, 1.0, font_size=44, color=WHITE)
hline(s, 1.12, 2.55, 17.5, color=GOOGLE_BLUE)

# Card Facilities
add_rect(s, 1.12, 2.8, 8.6, 7.8, fill=CARD_BG)
add_rect(s, 1.12, 2.8, 8.6, 0.1, fill=GOOGLE_BLUE)
add_tb(s, 'Facilities', 1.5, 3.05, 5.5, 0.65, font_size=22, bold=True, color=WHITE)
add_tb(s, 'R$ 2.200/mes', 7.0, 3.05, 2.4, 0.65, font_size=18, color=WHITE, align=PP_ALIGN.RIGHT)
kws_f = ['empresa de facilities SC','terceirizacao portaria Floripa',
         'gestao predial hospital','limpeza terceirizada universidade',
         'facilities Curitiba','zeladoria escola SC','facilities inst. financeira']
for i, kw in enumerate(kws_f):
    col = i % 2; row = i // 2
    add_rect(s, 1.4+col*4.1, 4.0+row*1.0, 3.8, 0.75, fill=RGBColor(0x14,0x32,0x38), radius=6000)
    add_tb(s, kw, 1.55+col*4.1, 4.08+row*1.0, 3.5, 0.55, font_size=13, color=RGBColor(0xA0,0xC8,0xFF))

# Card Seguranca
add_rect(s, 10.28, 2.8, 8.6, 7.8, fill=CARD_BG)
add_rect(s, 10.28, 2.8, 8.6, 0.1, fill=RGBColor(0x34,0xA8,0x53))
add_tb(s, 'Seguranca', 10.7, 3.05, 5.5, 0.65, font_size=22, bold=True, color=WHITE)
add_tb(s, 'R$ 1.800/mes', 16.2, 3.05, 2.4, 0.65, font_size=18, color=WHITE, align=PP_ALIGN.RIGHT)
kws_s = ['vigilancia patrimonial SC','empresa de seguranca Floripa',
         'seguranca hospitalar SC','vigilante terceirizado Curitiba',
         'seguranca escolar Joinville','seguranca inst. financeira']
for i, kw in enumerate(kws_s):
    col = i % 2; row = i // 2
    add_rect(s, 10.55+col*4.1, 4.0+row*1.0, 3.8, 0.75, fill=RGBColor(0x0F,0x28,0x1E), radius=6000)
    add_tb(s, kw, 10.7+col*4.1, 4.08+row*1.0, 3.5, 0.55, font_size=13, color=RGBColor(0x7D,0xCF,0xA0))

add_tb(s, 'Ad groups segmentados por nicho: copy e extensoes personalizadas para Saude, Ensino e Financeiro',
       1.12, 10.7, 17.5, 0.45, font_size=13, color=GRAY)

print('Slide 24 (Google Search) ok')

# ─────────────────────────────────────────
# SLIDE 25 — GOOGLE: DISPLAY (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

add_tb(s, 'Google Ads  |  R$ 1.000/mes  |  Display',
       1.12, 0.80, 12.0, 0.55, font_size=14, color=GOOGLE_BLUE)
add_tb(s, 'Alcancar decisores antes da busca ativa',
       1.12, 1.4, 14.0, 1.0, font_size=44, color=WHITE)
hline(s, 1.12, 2.55, 17.5, color=GOOGLE_BLUE)

blocos = [
    ('Segmentacao por cargo',
     'Diretores Administrativos, Gerentes de Facilities e Gestores Operacionais em SC e PR.',
     'In-Market'),
    ('Segmentacao por empresa',
     'Hospitais, universidades, cooperativas e bancos na rede Google + YouTube.',
     'Audience'),
    ('Alimenta o retargeting',
     'Quem clica no Display entra na LP e cai na audiencia de retargeting do Meta. Funil integrado.',
     'Funil'),
]
for i, (title, desc, tag) in enumerate(blocos):
    x = 1.12 + i * 6.2
    add_rect(s, x, 2.8, 5.8, 5.8, fill=CARD_BG)
    add_tb(s, title, x+0.4, 3.1, 5.0, 0.75, font_size=20, bold=True, color=WHITE)
    add_tb(s, desc, x+0.4, 4.0, 5.0, 3.5, font_size=16, color=GRAY)
    add_rect(s, x+0.4, 7.8, 1.6, 0.45, fill=RGBColor(0x14,0x32,0x46), radius=6000)
    add_tb(s, tag, x+0.4, 7.82, 1.6, 0.4, font_size=12, color=RGBColor(0xA0,0xC8,0xFF), align=PP_ALIGN.CENTER)

add_rect(s, 1.12, 9.0, 17.5, 1.55, fill=CARD_BG)
add_tb(s, 'Display e o gerador de audiencia para o retargeting - nao e onde o lead converte, e onde ele descobre a Orcali.',
       1.5, 9.25, 16.5, 1.0, font_size=16, color=WHITE)

print('Slide 25 (Google Display) ok')

# ─────────────────────────────────────────
# SLIDE 26 — LP + TYPEBOT (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

add_tb(s, 'Estrutura de Captacao',
       1.12, 0.80, 10.0, 0.55, font_size=14, color=GREEN)
add_tb(s, 'Landing Page + Formulario de qualificacao',
       1.12, 1.4, 16.0, 1.0, font_size=44, color=WHITE)
hline(s, 1.12, 2.55, 17.5)

# Painel LP
add_rect(s, 1.12, 2.8, 8.6, 7.8, fill=CARD_BG)
add_rect(s, 1.12, 2.8, 8.6, 0.85, fill=RGBColor(0x00,0x2E,0x6E))
add_tb(s, 'orcali.com.br/facilities', 1.5, 2.9, 5.5, 0.65, font_size=14, color=LIGHT_GRAY)
add_rect(s, 7.8, 2.88, 1.5, 0.52, fill=GREEN, radius=6000)
add_tb(s, 'PRONTO', 7.8, 2.9, 1.5, 0.48, font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
lp_items = [
    'Hero: "Operacoes que nao podem parar"',
    'Diferenciais: 58 anos  |  1.554 clientes',
    'Servicos por nicho: Saude, Ensino, Financeiro',
    'Prova social + certificacoes ISO',
    'Formulario integrado via Typebot',
    'LP Seguranca  ->  Mes 2',
]
for i, item in enumerate(lp_items):
    c = GRAY if i == 5 else WHITE
    add_tb(s, '->  ' + item, 1.4, 4.05+i*0.62, 8.0, 0.6, font_size=15, color=c)

# Painel Formulario
add_rect(s, 10.28, 2.8, 8.6, 7.8, fill=CARD_BG)
add_tb(s, 'Campos do formulario de qualificacao',
       10.7, 2.95, 7.5, 0.6, font_size=16, bold=True, color=GREEN)
fields = [
    'Nome completo',
    'Empresa',
    'CNPJ',
    'E-mail corporativo',
    'Telefone / WhatsApp',
    'Segmento  (Saude | Ensino | Financeiro)',
    'Cidade / Regiao',
]
for i, field in enumerate(fields):
    add_rect(s, 10.55, 3.7+i*0.62, 7.95, 0.52, fill=RGBColor(0x16,0x30,0x30), radius=4000)
    add_tb(s, '  ' + field, 10.75, 3.73+i*0.62, 7.6, 0.46, font_size=15, color=WHITE)

add_rect(s, 10.55, 8.15, 7.95, 0.7, fill=GREEN, radius=6000)
add_tb(s, 'Solicitar diagnostico operacional gratuito  ->',
       10.55, 8.17, 7.95, 0.66, font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

print('Slide 26 (LP + Typebot) ok')

# ─────────────────────────────────────────
# SLIDE 27 — INTEGRACOES (clone slide 6)
# ─────────────────────────────────────────
s = clone_slide(prs, 5)
remove_textboxes(s)

add_tb(s, 'Integracoes & Mensuracao',
       1.12, 0.80, 10.0, 0.55, font_size=14, color=GREEN)
add_tb(s, 'Do clique ao comercial rastreado',
       1.12, 1.4, 14.0, 1.0, font_size=44, color=WHITE)
hline(s, 1.12, 2.55, 17.5)

nodes = [
    ('Meta / Google',   'Pixel + CAPI\nGoogle Tag',   META_BLUE),
    ('Landing Page',    'Pixel fire\nViewContent',    GREEN),
    ('Planilha Google', 'Lead + UTMs\n+ origem',      RGBColor(0x0F,0x9D,0x58)),
    ('RD Station CRM',  'Nutricao\n+ funil',          RGBColor(0x8B,0x5C,0xF6)),
    ('Comercial',       'Lead pronto\np/ abordagem',  RGBColor(0xFB,0x92,0x3C)),
]
for i, (title, desc, color) in enumerate(nodes):
    x = 1.12 + i * 3.7
    add_rect(s, x, 3.0, 3.3, 4.2, fill=CARD_BG)
    add_rect(s, x, 3.0, 3.3, 0.1, fill=color)
    add_tb(s, title, x+0.2, 3.3, 2.9, 0.85, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(s, desc,  x+0.2, 4.25, 2.9, 1.4, font_size=15, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_tb(s, '>', x+3.3, 4.5, 0.4, 0.6, font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

labels = ['clique', 'form submit', 'Apps Script', 'alerta CRM']
for i, label in enumerate(labels):
    x = 1.12 + i * 3.7 + 3.3
    add_tb(s, label, x, 7.35, 0.7, 0.4, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s, 1.12, 8.1, 17.5, 2.45, fill=CARD_BG)
add_tb(s,
    'Cada lead chega com origem rastreada: plataforma | campanha | nicho | palavra-chave.',
    1.5, 8.35, 16.5, 0.7, font_size=17, bold=True, color=WHITE)
add_tb(s,
    'Resolve a principal dor do historico: R$ 6k/mes em 2024 sem mensuracao de resultados.',
    1.5, 9.1, 16.5, 0.7, font_size=16, color=GRAY)

print('Slide 27 (Integracoes) ok')

# ─────────────────────────────────────────
prs.save(DEST)
print(f'\nSalvo em: {DEST}')
print(f'Total de slides: {len(prs.slides)} (18 originais + 9 novos)')
